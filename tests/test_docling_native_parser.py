from __future__ import annotations

from io import BytesIO
from zipfile import ZIP_STORED, ZipFile

import pytest

pytest.importorskip("docling")

from docx import Document
from odfdo import Document as OdfDocument
from odfdo import DrawPage, Frame, Paragraph, Table
from openpyxl import Workbook
from pptx import Presentation

from grounded_docparse.config import ParserConfig
from grounded_docparse.docling_native import (
    make_docling_converter,
    make_docling_rapidocr_converter,
)
from grounded_docparse.native import (
    CellSourceAnchor,
    CsvSourceAnchor,
    ProcessingType,
    SourceFormat,
    StructuralSourceAnchor,
    TextSourceAnchor,
)
from grounded_docparse.native_parsers import DoclingNativeParser


def _parse(data: bytes, name: str, source_format: SourceFormat, processing_type: ProcessingType):
    return DoclingNativeParser(ParserConfig.from_env()).parse(
        data,
        name,
        source_format=source_format,
        processing_type=processing_type,
    )


def _docx() -> bytes:
    document = Document()
    document.add_paragraph("Exact paragraph")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Value"
    stream = BytesIO()
    document.save(stream)
    return stream.getvalue()


def _pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Exact shape"
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _multi_paragraph_pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title = slide.shapes.title.text_frame
    title.paragraphs[0].text = "Grounded DocParse"
    title.add_paragraph().text = "GitHub Overview"
    stream = BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def _xlsx() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Accounts"
    sheet["B2"] = "Name"
    sheet["C2"] = "Value"
    sheet["B3"] = "Ada"
    sheet["C3"] = "=40+2"
    stream = BytesIO()
    workbook.save(stream)
    workbook.close()
    return stream.getvalue()


def _epub() -> bytes:
    stream = BytesIO()
    with ZipFile(stream, "w") as archive:
        archive.writestr("mimetype", "application/epub+zip", compress_type=ZIP_STORED)
        archive.writestr(
            "META-INF/container.xml",
            '<container xmlns="urn:oasis:names:tc:opendocument:xmlns:container">'
            '<rootfiles><rootfile full-path="OEBPS/content.opf" '
            'media-type="application/oebps-package+xml"/></rootfiles></container>',
        )
        archive.writestr(
            "OEBPS/content.opf",
            '<package xmlns="http://www.idpf.org/2007/opf" version="3.0" '
            'unique-identifier="id"><metadata xmlns:dc="http://purl.org/dc/elements/1.1/">'
            '<dc:identifier id="id">fixture</dc:identifier><dc:title>Fixture</dc:title>'
            '<dc:language>en</dc:language></metadata><manifest>'
            '<item id="chapter" href="chapter.xhtml" media-type="application/xhtml+xml"/>'
            '</manifest><spine><itemref idref="chapter"/></spine></package>',
        )
        archive.writestr(
            "OEBPS/chapter.xhtml",
            '<html xmlns="http://www.w3.org/1999/xhtml"><body>'
            '<h1>Chapter title</h1><p>Chapter text</p></body></html>',
        )
    return stream.getvalue()


def _odf(kind: str) -> bytes:
    document = OdfDocument(kind)
    if kind == "text":
        document.body.append(Paragraph("ODT paragraph"))
    elif kind == "presentation":
        page = DrawPage(name="Named slide")
        frame = Frame(name="Text", presentation_class="body")
        frame.set_text_box("ODP shape")
        page.append(frame)
        document.body.append(page)
    else:
        table = Table("Data", width=2, height=1)
        table.set_value("A1", "Left")
        table.set_value("B1", "Right")
        document.body.append(table)
    stream = BytesIO()
    document.save(stream)
    return stream.getvalue()


def test_converter_allowlist_excludes_pdf_and_images_and_disables_models() -> None:
    from docling.datamodel.base_models import InputFormat
    from docling.pipeline.simple_pipeline import SimplePipeline

    converter = make_docling_converter()

    assert InputFormat.PDF not in converter.allowed_formats
    assert InputFormat.IMAGE not in converter.allowed_formats
    assert InputFormat.DOC not in converter.allowed_formats
    for option in converter.format_to_options.values():
        assert option.pipeline_cls is SimplePipeline
        assert option.pipeline_options.enable_remote_services is False
        assert option.pipeline_options.allow_external_plugins is False
        assert option.pipeline_options.do_picture_classification is False
        assert option.pipeline_options.do_picture_description is False
        assert option.pipeline_options.do_chart_extraction is False


def test_rapidocr_converter_is_image_only_cpu_and_offline(tmp_path) -> None:
    from docling.datamodel.accelerator_options import AcceleratorDevice
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import OcrMode, RapidOcrOptions

    converter = make_docling_rapidocr_converter(tmp_path)

    assert set(converter.allowed_formats) == {InputFormat.IMAGE}
    options = converter.format_to_options[InputFormat.IMAGE].pipeline_options
    assert options.artifacts_path == tmp_path
    assert options.accelerator_options.device is AcceleratorDevice.CPU
    assert options.do_ocr is True
    assert isinstance(options.ocr_options, RapidOcrOptions)
    assert options.ocr_options.backend == "onnxruntime"
    assert options.ocr_options.mode is OcrMode.LAYOUT_REGIONS
    assert options.layout_options.engine_options.compile_model is False
    assert options.enable_remote_services is False
    assert options.allow_external_plugins is False


def test_docx_maps_paragraph_table_and_cells_to_exact_paths() -> None:
    result = _parse(_docx(), "fixture.docx", SourceFormat.DOCX, ProcessingType.WORD)

    assert result.annotated_pdf is None
    assert result.document.base_text == "Exact paragraph\nName\tValue"
    assert isinstance(result.document.elements[0].source.anchor, StructuralSourceAnchor)
    assert result.document.elements[0].source.anchor.path == "/word/document/body/p[1]"
    table = next(item for item in result.document.elements if item.type == "table")
    assert table.source.anchor.path == "/word/document/body/tbl[1]"
    assert [
        result.document.elements[int(child.split("-")[1]) - 1].source.anchor.path
        for child in table.children
    ] == [
        "/word/document/body/tbl[1]/row[1]/cell[1]",
        "/word/document/body/tbl[1]/row[1]/cell[2]",
    ]


def test_pptx_maps_text_to_slide_and_shape_id() -> None:
    result = _parse(
        _pptx(), "fixture.pptx", SourceFormat.PPTX, ProcessingType.POWERPOINT
    )

    anchor = result.document.elements[0].source.anchor
    assert result.document.units[0].kind == "slide"
    assert anchor.unit_id == "slide-1"
    assert anchor.path.startswith("/ppt/slides/slide[1]/shape[id=")


def test_pptx_maps_each_shape_paragraph_to_its_exact_source_path() -> None:
    result = _parse(
        _multi_paragraph_pptx(),
        "fixture.pptx",
        SourceFormat.PPTX,
        ProcessingType.POWERPOINT,
    )

    elements = {
        item.text: item
        for item in result.document.elements
        if item.text in {"Grounded DocParse", "GitHub Overview"}
    }
    assert set(elements) == {"Grounded DocParse", "GitHub Overview"}
    first_path = elements["Grounded DocParse"].source.anchor.path
    second_path = elements["GitHub Overview"].source.anchor.path
    assert first_path.endswith("/paragraph[1]")
    assert second_path.endswith("/paragraph[2]")
    assert first_path.rsplit("/paragraph", 1)[0] == second_path.rsplit("/paragraph", 1)[0]


def test_xlsx_preserves_formula_and_maps_each_cell() -> None:
    result = _parse(_xlsx(), "fixture.xlsx", SourceFormat.XLSX, ProcessingType.EXCEL)

    assert "=40+2" in result.document.base_text
    anchors = {
        item.text: item.source.anchor.cell_range
        for item in result.document.elements
        if item.type == "table_cell"
    }
    assert anchors == {"Name": "B2", "Value": "C2", "Ada": "B3", "=40+2": "C3"}
    assert all(isinstance(item.source.anchor, CellSourceAnchor) for item in result.document.elements)


def test_csv_maps_cells_to_one_based_rows_and_columns() -> None:
    result = _parse(
        b'Name,Note\nAda,"line one\nline two"\n',
        "fixture.csv",
        SourceFormat.CSV,
        ProcessingType.CSV,
    )

    anchors = [
        item.source.anchor
        for item in result.document.elements
        if item.type == "table_cell"
    ]
    assert all(isinstance(anchor, CsvSourceAnchor) for anchor in anchors)
    assert (anchors[-1].row_start, anchors[-1].column_start) == (2, 2)


def test_html_records_remote_image_without_fetching_or_ocr() -> None:
    result = _parse(
        b'<html><body><h1>Title</h1><p>Text</p><img src="https://example.invalid/x.png" alt="Chart"></body></html>',
        "fixture.html",
        SourceFormat.HTML,
        ProcessingType.OTHER_NATIVE,
    )

    assert [item.source.anchor.path for item in result.document.elements] == [
        "/html[1]/body[1]/h1[1]",
        "/html[1]/body[1]/p[1]",
    ]
    assert result.document.assets[0].reference == "https://example.invalid/x.png"
    assert result.document.assets[0].sha256 is None
    assert result.document.assets[0].ocr_performed is False


def test_epub_uses_spine_item_and_dom_path() -> None:
    result = _parse(
        _epub(), "fixture.epub", SourceFormat.EPUB, ProcessingType.OTHER_NATIVE
    )

    assert result.document.units[0].label == "chapter.xhtml"
    assert result.document.elements[0].source.anchor.path.startswith(
        "epub:chapter.xhtml:/html[1]/body[1]/h1[1]"
    )


def test_markdown_uses_exact_line_and_column_anchor() -> None:
    result = _parse(
        b"# Title\n\nExact paragraph\n",
        "fixture.md",
        SourceFormat.MARKDOWN,
        ProcessingType.OTHER_NATIVE,
    )

    anchors = [item.source.anchor for item in result.document.elements]
    assert all(isinstance(anchor, TextSourceAnchor) for anchor in anchors)
    assert (anchors[1].start_line, anchors[1].end_line) == (3, 3)


@pytest.mark.parametrize(
    ("kind", "source_format", "expected", "unit_kind"),
    [
        ("text", SourceFormat.ODT, "ODT paragraph", "document"),
        ("presentation", SourceFormat.ODP, "ODP shape", "slide"),
        ("spreadsheet", SourceFormat.ODS, "Left\tRight", "sheet"),
    ],
)
def test_open_document_formats_are_native_and_exactly_anchored(
    kind: str, source_format: SourceFormat, expected: str, unit_kind: str
) -> None:
    result = _parse(
        _odf(kind),
        f"fixture.{source_format.value}",
        source_format,
        ProcessingType.OTHER_NATIVE,
    )

    assert expected in result.document.base_text
    assert result.document.units[0].kind == unit_kind
    assert all(item.source.anchor.unit_id in {unit.id for unit in result.document.units} for item in result.document.elements)


def test_every_textual_element_span_points_into_immutable_base_text() -> None:
    result = _parse(_docx(), "fixture.docx", SourceFormat.DOCX, ProcessingType.WORD)

    for element in result.document.elements:
        assert (
            result.document.base_text[element.source.start : element.source.end]
            == element.text
        )
